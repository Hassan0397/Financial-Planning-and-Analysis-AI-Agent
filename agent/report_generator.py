"""
Report Generator Module
Professional FP&A report generation with PDF, Word, PowerPoint outputs
"""

import streamlit as st
import pandas as pd
import numpy as np
from typing import Dict, List, Tuple, Optional, Any, Union
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime, timedelta
import os
import tempfile
import base64
import json
import warnings
warnings.filterwarnings('ignore')

# Try to import reporting libraries
try:
    from fpdf import FPDF
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    reportlab_available = True
except ImportError:
    reportlab_available = False
    st.warning("⚠️ ReportLab not installed. Install with: pip install reportlab")

try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    python_docx_available = True
except ImportError:
    python_docx_available = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    python_pptx_available = True
except ImportError:
    python_pptx_available = False


class ReportGenerator:
    """
    Professional report generation for FP&A analysis
    """
    
    def __init__(self):
        self.report_templates = self._initialize_report_templates()
        self.company_info = self._get_default_company_info()
        self.report_sections = self._initialize_report_sections()
    
    def _initialize_report_templates(self) -> Dict[str, Dict]:
        """Initialize report templates"""
        return {
            'executive_summary': {
                'name': 'Executive Summary',
                'description': 'High-level overview for executives',
                'sections': ['executive_overview', 'key_findings', 'recommendations'],
                'format': 'concise',
                'target_audience': 'C-level executives'
            },
            'monthly_performance': {
                'name': 'Monthly Performance Report',
                'description': 'Detailed monthly performance analysis',
                'sections': ['performance_overview', 'revenue_analysis', 'expense_analysis', 
                           'profitability', 'cash_flow', 'forecast', 'scenarios'],
                'format': 'detailed',
                'target_audience': 'Department heads'
            },
            'quarterly_review': {
                'name': 'Quarterly Business Review',
                'description': 'Comprehensive quarterly review',
                'sections': ['executive_summary', 'financial_performance', 'operational_metrics',
                           'market_analysis', 'forecast_outlook', 'risk_assessment', 'action_plan'],
                'format': 'comprehensive',
                'target_audience': 'Board of Directors'
            },
            'budget_variance': {
                'name': 'Budget vs Actual Analysis',
                'description': 'Variance analysis report',
                'sections': ['variance_summary', 'revenue_variance', 'expense_variance',
                           'profit_variance', 'root_cause_analysis', 'corrective_actions'],
                'format': 'analytical',
                'target_audience': 'Finance team'
            },
            'forecast_report': {
                'name': 'Financial Forecast Report',
                'description': 'Detailed forecast analysis',
                'sections': ['forecast_summary', 'revenue_forecast', 'expense_forecast',
                           'cash_flow_forecast', 'scenario_analysis', 'key_assumptions',
                           'sensitivity_analysis'],
                'format': 'forward_looking',
                'target_audience': 'Planning committee'
            },
            'ad_hoc_analysis': {
                'name': 'Ad-hoc Analysis Report',
                'description': 'Custom analysis report',
                'sections': ['custom'],
                'format': 'flexible',
                'target_audience': 'Custom'
            }
        }
    
    def _get_default_company_info(self) -> Dict[str, str]:
        """Get default company information"""
        return {
            'company_name': 'Your Company',
            'logo_path': 'assets/logo.png' if os.path.exists('assets/logo.png') else None,
            'address': '123 Business Street, Suite 100',
            'city_state_zip': 'New York, NY 10001',
            'phone': '(123) 456-7890',
            'email': 'info@yourcompany.com',
            'website': 'www.yourcompany.com',
            'fiscal_year': datetime.now().year,
            'report_period': 'Q1 2024',
            'prepared_by': 'FP&A Department',
            'prepared_date': datetime.now().strftime('%B %d, %Y')
        }
    
    def _initialize_report_sections(self) -> Dict[str, Dict]:
        """Initialize report section templates"""
        return {
            'executive_overview': {
                'title': 'Executive Overview',
                'description': 'High-level summary of financial performance',
                'content_type': 'text',
                'default_text': 'This section provides a high-level overview of the company\'s financial performance during the reporting period.'
            },
            'key_findings': {
                'title': 'Key Findings',
                'description': 'Most important insights from the analysis',
                'content_type': 'bullet_points',
                'default_text': 'Key insights and discoveries from the financial analysis.'
            },
            'recommendations': {
                'title': 'Recommendations',
                'description': 'Actionable recommendations based on analysis',
                'content_type': 'bullet_points',
                'default_text': 'Strategic recommendations for management consideration.'
            },
            'performance_overview': {
                'title': 'Performance Overview',
                'description': 'Overall financial performance metrics',
                'content_type': 'kpi_metrics',
                'default_text': 'Summary of key performance indicators.'
            },
            'revenue_analysis': {
                'title': 'Revenue Analysis',
                'description': 'Detailed revenue breakdown and analysis',
                'content_type': 'charts_tables',
                'default_text': 'Analysis of revenue streams, trends, and drivers.'
            },
            'expense_analysis': {
                'title': 'Expense Analysis',
                'description': 'Detailed expense breakdown and analysis',
                'content_type': 'charts_tables',
                'default_text': 'Analysis of expense categories and cost drivers.'
            },
            'profitability': {
                'title': 'Profitability Analysis',
                'description': 'Profit margins and profitability metrics',
                'content_type': 'charts_tables',
                'default_text': 'Analysis of profitability across products, segments, and periods.'
            },
            'cash_flow': {
                'title': 'Cash Flow Analysis',
                'description': 'Cash flow statement and analysis',
                'content_type': 'charts_tables',
                'default_text': 'Analysis of cash flow from operations, investing, and financing activities.'
            },
            'forecast': {
                'title': 'Forecast Outlook',
                'description': 'Future projections and forecasts',
                'content_type': 'forecast_charts',
                'default_text': 'Financial forecasts and projections for upcoming periods.'
            },
            'scenarios': {
                'title': 'Scenario Analysis',
                'description': 'What-if scenario analysis',
                'content_type': 'scenario_tables',
                'default_text': 'Analysis of different business scenarios and their financial impact.'
            },
            'variance_summary': {
                'title': 'Variance Summary',
                'description': 'Summary of budget vs actual variances',
                'content_type': 'variance_tables',
                'default_text': 'Summary of significant variances between budgeted and actual results.'
            },
            'root_cause_analysis': {
                'title': 'Root Cause Analysis',
                'description': 'Analysis of variance root causes',
                'content_type': 'text',
                'default_text': 'Investigation into the underlying causes of significant variances.'
            },
            'corrective_actions': {
                'title': 'Corrective Actions',
                'description': 'Proposed actions to address variances',
                'content_type': 'bullet_points',
                'default_text': 'Recommended actions to address identified issues and variances.'
            }
        }
    
    def display_report_interface(self, dataframes: Dict[str, pd.DataFrame]) -> Dict[str, Any]:
        """
        Main Streamlit interface for report generation
        """
        st.subheader("📄 Professional Report Generation")
        
        # Initialize session state for report
        if 'report_content' not in st.session_state:
            st.session_state.report_content = {}
        
        # Check for available data
        has_data = bool(dataframes)
        has_analytics = 'analysis_results' in st.session_state and st.session_state.analysis_results
        has_forecasts = 'forecast_results' in st.session_state and st.session_state.forecast_results
        has_scenarios = 'scenario_results' in st.session_state and st.session_state.scenario_results
        has_visualizations = 'viz_results' in st.session_state and st.session_state.viz_results
        
        # Create tabs for different report sections
        tab1, tab2, tab3, tab4, tab5 = st.tabs([
            "📋 Report Setup", 
            "📊 Content Selection", 
            "🎨 Formatting",
            "📤 Generate Report",
            "📁 Report Library"
        ])
        
        report_results = {}
        
        with tab1:
            report_results['setup'] = self._report_setup_tab(has_data, has_analytics, has_forecasts, has_scenarios)
        
        with tab2:
            report_results['content'] = self._content_selection_tab()
        
        with tab3:
            report_results['formatting'] = self._formatting_tab()
        
        with tab4:
            report_results['generation'] = self._generate_report_tab()
        
        with tab5:
            report_results['library'] = self._report_library_tab()
        
        return report_results
    
    def _report_setup_tab(self, has_data: bool, has_analytics: bool, 
                         has_forecasts: bool, has_scenarios: bool) -> Dict[str, Any]:
        """
        Report setup and template selection
        """
        st.markdown("### 📋 Report Setup")
        
        results = {}
        
        # Company information
        st.markdown("#### 🏢 Company Information")
        
        col1, col2 = st.columns(2)
        
        with col1:
            self.company_info['company_name'] = st.text_input(
                "Company Name:",
                value=self.company_info['company_name']
            )
            
            self.company_info['report_period'] = st.text_input(
                "Report Period:",
                value=self.company_info['report_period'],
                placeholder="e.g., Q1 2024, January 2024, FY2023"
            )
        
        with col2:
            self.company_info['prepared_by'] = st.text_input(
                "Prepared By:",
                value=self.company_info['prepared_by']
            )
            
            self.company_info['prepared_date'] = st.text_input(
                "Report Date:",
                value=self.company_info['prepared_date']
            )
        
        # Report template selection
        st.markdown("#### 📄 Report Template")
        
        template_options = list(self.report_templates.keys())
        template_names = [self.report_templates[t]['name'] for t in template_options]
        
        selected_template = st.selectbox(
            "Select report template:",
            options=template_options,
            format_func=lambda x: self.report_templates[x]['name'],
            help="Choose a pre-built template or create a custom report"
        )
        
        if selected_template:
            template = self.report_templates[selected_template]
            
            st.markdown(f"**Template:** {template['name']}")
            st.markdown(f"**Description:** {template['description']}")
            st.markdown(f"**Target Audience:** {template['target_audience']}")
            
            # Show available data status
            st.markdown("#### 📊 Available Data Status")
            
            status_cols = st.columns(4)
            
            with status_cols[0]:
                st.metric(
                    "Data Files",
                    "✅ Available" if has_data else "❌ Missing",
                    delta_color="normal" if has_data else "inverse"
                )
            
            with status_cols[1]:
                st.metric(
                    "Analytics",
                    "✅ Available" if has_analytics else "⚠️ Recommended",
                    delta_color="normal" if has_analytics else "off"
                )
            
            with status_cols[2]:
                st.metric(
                    "Forecasts",
                    "✅ Available" if has_forecasts else "⚪ Optional",
                    delta_color="normal" if has_forecasts else "off"
                )
            
            with status_cols[3]:
                st.metric(
                    "Scenarios",
                    "✅ Available" if has_scenarios else "⚪ Optional",
                    delta_color="normal" if has_scenarios else "off"
                )
            
            # Data source selection
            if has_data:
                st.markdown("#### 📁 Data Sources")
                
                data_sources = st.multiselect(
                    "Select data sources to include:",
                    options=list(st.session_state.get('uploaded_data', {}).keys()),
                    default=list(st.session_state.get('uploaded_data', {}).keys())[:1] if st.session_state.get('uploaded_data') else [],
                    help="Select which data files to include in the report"
                )
                
                results['data_sources'] = data_sources
            
            # Template customization
            st.markdown("#### ⚙️ Template Customization")
            
            custom_sections = st.multiselect(
                "Customize report sections:",
                options=list(self.report_sections.keys()),
                default=template['sections'] if selected_template != 'ad_hoc_analysis' else [],
                help="Select which sections to include in the report"
            )
            
            results['template'] = selected_template
            results['sections'] = custom_sections
            results['company_info'] = self.company_info.copy()
            
            # Save report configuration
            if st.button("💾 Save Report Configuration", use_container_width=True):
                st.session_state.report_config = results
                st.success("Report configuration saved!")
        
        return results
    
    def _content_selection_tab(self) -> Dict[str, Any]:
        """
        Content selection and customization
        """
        st.markdown("### 📊 Report Content")
        
        results = {}
        
        # Check if report is configured
        if 'report_config' not in st.session_state:
            st.info("Please set up the report first in the 'Report Setup' tab.")
            return results
        
        config = st.session_state.report_config
        
        # Content selection for each section
        st.markdown(f"#### Configuring: {self.report_templates[config['template']]['name']}")
        
        for section_key in config.get('sections', []):
            section = self.report_sections[section_key]
            
            with st.expander(f"📝 {section['title']}", expanded=True):
                st.markdown(f"**Description:** {section['description']}")
                
                # Content type handling
                if section['content_type'] == 'text':
                    content = st.text_area(
                        f"Content for {section['title']}:",
                        value=section['default_text'],
                        height=150,
                        key=f"text_{section_key}"
                    )
                    results[section_key] = {'type': 'text', 'content': content}
                
                elif section['content_type'] == 'bullet_points':
                    st.markdown("**Add bullet points:**")
                    
                    bullet_points = []
                    for i in range(5):
                        point = st.text_input(
                            f"Point {i+1}:",
                            value="",
                            key=f"bullet_{section_key}_{i}",
                            placeholder=f"Enter point {i+1} (leave empty to skip)"
                        )
                        if point:
                            bullet_points.append(point)
                    
                    # Add more button
                    if st.button(f"➕ Add more points to {section['title']}", key=f"add_{section_key}"):
                        # This would add more input fields in a real implementation
                        st.info("In a full implementation, this would add more input fields")
                    
                    results[section_key] = {'type': 'bullet_points', 'content': bullet_points}
                
                elif section['content_type'] == 'kpi_metrics':
                    st.markdown("**Select KPIs to include:**")
                    
                    # Get available KPIs from analytics
                    available_kpis = []
                    if 'analysis_results' in st.session_state:
                        analysis = st.session_state.analysis_results
                        if 'kpi_dashboard' in analysis:
                            kpi_data = analysis['kpi_dashboard']
                            available_kpis = list(kpi_data.keys())
                    
                    if available_kpis:
                        selected_kpis = st.multiselect(
                            f"Select KPIs for {section['title']}:",
                            options=available_kpis,
                            default=available_kpis[:min(8, len(available_kpis))],
                            key=f"kpis_{section_key}"
                        )
                        
                        # Format options
                        kpi_format = st.selectbox(
                            f"Format for {section['title']}:",
                            options=['Table', 'Cards', 'Chart'],
                            key=f"kpi_format_{section_key}"
                        )
                        
                        results[section_key] = {
                            'type': 'kpi_metrics',
                            'kpis': selected_kpis,
                            'format': kpi_format
                        }
                    else:
                        st.warning("No KPI data available. Please run analytics first.")
                        results[section_key] = {'type': 'kpi_metrics', 'kpis': [], 'format': 'Table'}
                
                elif section['content_type'] == 'charts_tables':
                    st.markdown("**Select charts and tables:**")
                    
                    # Get available visualizations
                    available_viz = []
                    if 'viz_results' in st.session_state:
                        viz = st.session_state.viz_results
                        # Extract chart information
                        for tab_name, tab_content in viz.items():
                            if 'chart' in tab_content:
                                available_viz.append(f"{tab_name}_chart")
                    
                    if available_viz:
                        selected_viz = st.multiselect(
                            f"Select visualizations for {section['title']}:",
                            options=available_viz,
                            default=available_viz[:min(3, len(available_viz))],
                            key=f"viz_{section_key}"
                        )
                        
                        results[section_key] = {
                            'type': 'charts_tables',
                            'visualizations': selected_viz
                        }
                    else:
                        st.warning("No visualization data available. Please create visualizations first.")
                        results[section_key] = {'type': 'charts_tables', 'visualizations': []}
                
                elif section['content_type'] == 'forecast_charts':
                    st.markdown("**Select forecast charts:**")
                    
                    if 'forecast_results' in st.session_state:
                        forecast_data = st.session_state.forecast_results
                        
                        forecast_options = []
                        if 'single_series' in forecast_data:
                            forecast_options.append('Single Series Forecast')
                        if 'scenarios' in forecast_data:
                            forecast_options.append('Forecast Scenarios')
                        if 'model_comparison' in forecast_data:
                            forecast_options.append('Model Comparison')
                        
                        selected_forecasts = st.multiselect(
                            f"Select forecasts for {section['title']}:",
                            options=forecast_options,
                            default=forecast_options,
                            key=f"forecasts_{section_key}"
                        )
                        
                        results[section_key] = {
                            'type': 'forecast_charts',
                            'forecasts': selected_forecasts
                        }
                    else:
                        st.warning("No forecast data available. Please run forecasting first.")
                        results[section_key] = {'type': 'forecast_charts', 'forecasts': []}
                
                elif section['content_type'] == 'scenario_tables':
                    st.markdown("**Select scenario analyses:**")
                    
                    if 'scenario_results' in st.session_state:
                        scenario_data = st.session_state.scenario_results
                        
                        scenario_options = []
                        if 'created_scenarios' in scenario_data:
                            scenario_options.append('Created Scenarios')
                        if 'comparison' in scenario_data:
                            scenario_options.append('Scenario Comparison')
                        if 'impact_analysis' in scenario_data:
                            scenario_options.append('Impact Analysis')
                        
                        selected_scenarios = st.multiselect(
                            f"Select scenarios for {section['title']}:",
                            options=scenario_options,
                            default=scenario_options,
                            key=f"scenarios_{section_key}"
                        )
                        
                        results[section_key] = {
                            'type': 'scenario_tables',
                            'scenarios': selected_scenarios
                        }
                    else:
                        st.warning("No scenario data available. Please run scenario analysis first.")
                        results[section_key] = {'type': 'scenario_tables', 'scenarios': []}
                
                elif section['content_type'] == 'variance_tables':
                    st.markdown("**Select variance analyses:**")
                    
                    # Check for variance data in analytics
                    if 'analysis_results' in st.session_state:
                        analysis = st.session_state.analysis_results
                        has_variance = 'variance_analysis' in analysis
                        
                        if has_variance:
                            variance_type = st.selectbox(
                                f"Variance analysis type for {section['title']}:",
                                options=['Summary Table', 'Detailed Breakdown', 'Root Cause Analysis'],
                                key=f"variance_type_{section_key}"
                            )
                            
                            results[section_key] = {
                                'type': 'variance_tables',
                                'variance_type': variance_type
                            }
                        else:
                            st.warning("No variance analysis available. Please run analytics with scenario data.")
                            results[section_key] = {'type': 'variance_tables', 'variance_type': 'Summary Table'}
                    else:
                        st.warning("No analytics data available. Please run analytics first.")
                        results[section_key] = {'type': 'variance_tables', 'variance_type': 'Summary Table'}
        
        # Save content configuration
        if st.button("💾 Save Content Configuration", type="primary", use_container_width=True):
            st.session_state.report_content = results
            st.success("Report content saved!")
        
        return results
    
    def _formatting_tab(self) -> Dict[str, Any]:
        """
        Report formatting and styling
        """
        st.markdown("### 🎨 Report Formatting")
        
        results = {}
        
        # Format options
        st.markdown("#### 📄 Format Options")
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            report_format = st.selectbox(
                "Output Format:",
                options=['PDF', 'Word (DOCX)', 'PowerPoint (PPTX)', 'HTML', 'All Formats'],
                index=0
            )
        
        with col2:
            paper_size = st.selectbox(
                "Paper Size:",
                options=['Letter', 'A4', 'Legal', 'Executive'],
                index=0
            )
        
        with col3:
            orientation = st.selectbox(
                "Orientation:",
                options=['Portrait', 'Landscape'],
                index=0
            )
        
        # Color scheme
        st.markdown("#### 🎨 Color Scheme")
        
        color_scheme = st.selectbox(
            "Color Scheme:",
            options=['Corporate Blue', 'Professional Gray', 'Green Finance', 'Custom']
        )
        
        if color_scheme == 'Custom':
            col1, col2, col3 = st.columns(3)
            with col1:
                primary_color = st.color_picker("Primary Color", value="#2E86AB")
            with col2:
                secondary_color = st.color_picker("Secondary Color", value="#A23B72")
            with col3:
                accent_color = st.color_picker("Accent Color", value="#F18F01")
            
            colors = {
                'primary': primary_color,
                'secondary': secondary_color,
                'accent': accent_color
            }
        else:
            color_palettes = {
                'Corporate Blue': {'primary': '#2E86AB', 'secondary': '#6B9AC4', 'accent': '#97D8C4'},
                'Professional Gray': {'primary': '#4A5568', 'secondary': '#718096', 'accent': '#A0AEC0'},
                'Green Finance': {'primary': '#2D936C', 'secondary': '#5BBF95', 'accent': '#94D7BD'}
            }
            colors = color_palettes[color_scheme]
        
        # Font selection
        st.markdown("#### 🔤 Typography")
        
        font_family = st.selectbox(
            "Font Family:",
            options=['Arial', 'Times New Roman', 'Calibri', 'Helvetica', 'Georgia']
        )
        
        col1, col2 = st.columns(2)
        with col1:
            base_font_size = st.slider("Base Font Size:", 10, 14, 11)
        with col2:
            line_spacing = st.slider("Line Spacing:", 1.0, 2.0, 1.5, 0.1)
        
        # Header and footer
        st.markdown("#### 📝 Header & Footer")
        
        show_header = st.checkbox("Show header", value=True)
        show_footer = st.checkbox("Show footer", value=True)
        
        if show_header:
            header_text = st.text_input("Header text:", value="CONFIDENTIAL - INTERNAL USE ONLY")
        
        if show_footer:
            footer_text = st.text_input("Footer text:", value="Page {page_number} of {total_pages}")
        
        # Page numbering
        page_numbering = st.selectbox(
            "Page Numbering:",
            options=['Bottom Center', 'Bottom Right', 'Top Right', 'None']
        )
        
        # Watermark
        watermark = st.checkbox("Add watermark", value=False)
        if watermark:
            watermark_text = st.text_input("Watermark text:", value="DRAFT")
            watermark_opacity = st.slider("Watermark opacity:", 0.1, 0.5, 0.2, 0.05)
        
        # Save formatting
        formatting_config = {
            'report_format': report_format,
            'paper_size': paper_size,
            'orientation': orientation,
            'colors': colors,
            'font_family': font_family,
            'base_font_size': base_font_size,
            'line_spacing': line_spacing,
            'show_header': show_header,
            'show_footer': show_footer,
            'header_text': header_text if show_header else '',
            'footer_text': footer_text if show_footer else '',
            'page_numbering': page_numbering,
            'watermark': watermark,
            'watermark_text': watermark_text if watermark else '',
            'watermark_opacity': watermark_opacity if watermark else 0.2
        }
        
        if st.button("💾 Save Formatting", type="primary", use_container_width=True):
            st.session_state.report_formatting = formatting_config
            st.success("Formatting saved!")
        
        results = formatting_config
        return results
    
    def _generate_report_tab(self) -> Dict[str, Any]:
        """
        Generate and export report
        """
        st.markdown("### 📤 Generate Report")
        
        results = {}
        
        # Check if report is ready
        config_ready = 'report_config' in st.session_state
        content_ready = 'report_content' in st.session_state and st.session_state.report_content
        formatting_ready = 'report_formatting' in st.session_state
        
        status_cols = st.columns(3)
        
        with status_cols[0]:
            st.metric(
                "Setup",
                "✅ Complete" if config_ready else "❌ Required",
                delta_color="normal" if config_ready else "inverse"
            )
        
        with status_cols[1]:
            st.metric(
                "Content",
                "✅ Complete" if content_ready else "⚠️ Recommended",
                delta_color="normal" if content_ready else "off"
            )
        
        with status_cols[2]:
            st.metric(
                "Formatting",
                "✅ Complete" if formatting_ready else "⚪ Optional",
                delta_color="normal" if formatting_ready else "off"
            )
        
        if not config_ready:
            st.warning("Please complete report setup first.")
            return results
        
        # Report preview
        st.markdown("#### 👁️ Report Preview")
        
        # Generate preview
        if st.button("🔍 Generate Preview", use_container_width=True):
            with st.spinner("Generating report preview..."):
                preview = self._generate_report_preview()
                
                if preview:
                    st.markdown("##### 📋 Report Structure")
                    
                    # Show table of contents
                    toc_data = []
                    if 'report_config' in st.session_state:
                        config = st.session_state.report_config
                        for i, section_key in enumerate(config.get('sections', []), 1):
                            section = self.report_sections[section_key]
                            toc_data.append({
                                'Section': f"{i}. {section['title']}",
                                'Type': section['content_type'],
                                'Status': '✅' if section_key in st.session_state.report_content else '⚠️'
                            })
                    
                    toc_df = pd.DataFrame(toc_data)
                    st.dataframe(toc_df, use_container_width=True)
                    
                    # Show sample content
                    st.markdown("##### 📝 Sample Content")
                    
                    if 'report_content' in st.session_state:
                        content = st.session_state.report_content
                        for section_key, section_content in list(content.items())[:3]:  # Show first 3 sections
                            section = self.report_sections.get(section_key, {'title': section_key})
                            with st.expander(f"{section['title']} Preview"):
                                if section_content['type'] == 'text':
                                    st.write(section_content['content'][:500] + "...")
                                elif section_content['type'] == 'bullet_points':
                                    for point in section_content['content'][:3]:  # Show first 3 points
                                        st.write(f"• {point}")
                                elif section_content['type'] == 'kpi_metrics':
                                    st.write(f"KPIs: {len(section_content.get('kpis', []))} metrics selected")
        
        # Report generation options
        st.markdown("---")
        st.markdown("#### 🚀 Generate Final Report")
        
        report_name = st.text_input(
            "Report Name:",
            value=f"FP&A_Report_{datetime.now().strftime('%Y%m%d')}",
            help="Name for the generated report file"
        )
        
        include_appendix = st.checkbox("Include appendix with raw data", value=False)
        include_audit_trail = st.checkbox("Include audit trail", value=True)
        compress_files = st.checkbox("Compress output files", value=True)
        
        # Generate buttons for different formats
        st.markdown("##### 📄 Generate in Format:")
        
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            if st.button("📊 PDF Report", use_container_width=True):
                if reportlab_available:
                    with st.spinner("Generating PDF report..."):
                        pdf_report = self._generate_pdf_report()
                        if pdf_report:
                            results['pdf'] = pdf_report
                            st.success("PDF report generated successfully!")
                            
                            # Provide download link
                            self._provide_download_link(pdf_report, f"{report_name}.pdf", "application/pdf")
                else:
                    st.error("ReportLab not available. Install with: pip install reportlab")
        
        with col2:
            if st.button("📝 Word Document", use_container_width=True):
                if python_docx_available:
                    with st.spinner("Generating Word document..."):
                        docx_report = self._generate_docx_report()
                        if docx_report:
                            results['docx'] = docx_report
                            st.success("Word document generated successfully!")
                            
                            # Provide download link
                            self._provide_download_link(docx_report, f"{report_name}.docx", 
                                                       "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                else:
                    st.error("python-docx not available. Install with: pip install python-docx")
        
        with col3:
            if st.button("📈 PowerPoint", use_container_width=True):
                if python_pptx_available:
                    with st.spinner("Generating PowerPoint presentation..."):
                        pptx_report = self._generate_pptx_report()
                        if pptx_report:
                            results['pptx'] = pptx_report
                            st.success("PowerPoint presentation generated successfully!")
                            
                            # Provide download link
                            self._provide_download_link(pptx_report, f"{report_name}.pptx",
                                                       "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                else:
                    st.error("python-pptx not available. Install with: pip install python-pptx")
        
        with col4:
            if st.button("🌐 HTML Report", use_container_width=True):
                with st.spinner("Generating HTML report..."):
                    html_report = self._generate_html_report()
                    if html_report:
                        results['html'] = html_report
                        st.success("HTML report generated successfully!")
                        
                        # Provide download link
                        self._provide_download_link(html_report, f"{report_name}.html", "text/html")
        
        # Generate all formats
        st.markdown("---")
        if st.button("🚀 Generate All Formats", type="primary", use_container_width=True):
            with st.spinner("Generating reports in all formats..."):
                all_reports = self._generate_all_formats()
                
                if all_reports:
                    results['all'] = all_reports
                    st.success("All report formats generated successfully!")
                    
                    # Create zip file if multiple formats
                    if len(all_reports) > 1 and compress_files:
                        zip_file = self._create_zip_file(all_reports, report_name)
                        if zip_file:
                            self._provide_download_link(zip_file, f"{report_name}.zip", "application/zip")
        
        return results
    
    def _generate_report_preview(self) -> Dict[str, Any]:
        """Generate report preview"""
        preview = {
            'status': 'preview',
            'generated_at': datetime.now().isoformat(),
            'config': st.session_state.get('report_config', {}),
            'content_summary': {},
            'estimated_pages': 0
        }
        
        # Estimate pages based on content
        if 'report_content' in st.session_state:
            content = st.session_state.report_content
            preview['content_summary'] = {
                'total_sections': len(content),
                'section_types': {},
                'content_volume': 0
            }
            
            for section_key, section_content in content.items():
                content_type = section_content['type']
                if content_type not in preview['content_summary']['section_types']:
                    preview['content_summary']['section_types'][content_type] = 0
                preview['content_summary']['section_types'][content_type] += 1
                
                # Estimate content volume
                if content_type == 'text':
                    preview['content_summary']['content_volume'] += len(section_content.get('content', '')) / 1000
                elif content_type == 'bullet_points':
                    preview['content_summary']['content_volume'] += len(section_content.get('content', [])) * 0.5
                elif content_type == 'kpi_metrics':
                    preview['content_summary']['content_volume'] += len(section_content.get('kpis', [])) * 0.3
        
        # Estimate pages (rough estimate)
        content_volume = preview['content_summary'].get('content_volume', 0)
        preview['estimated_pages'] = max(1, int(content_volume * 0.5))
        
        return preview
    
    def _generate_pdf_report(self) -> Optional[bytes]:
        """Generate PDF report"""
        try:
            # Create PDF
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.lib import colors
            
            # Get formatting
            formatting = st.session_state.get('report_formatting', {})
            paper_size = formatting.get('paper_size', 'Letter')
            orientation = formatting.get('orientation', 'Portrait')
            
            # Create buffer for PDF
            import io
            buffer = io.BytesIO()
            
            # Set page size
            if paper_size == 'A4':
                pagesize = A4
            else:
                pagesize = letter
            
            if orientation == 'Landscape':
                pagesize = (pagesize[1], pagesize[0])
            
            # Create document
            doc = SimpleDocTemplate(
                buffer,
                pagesize=pagesize,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            # Build story (content)
            story = []
            styles = getSampleStyleSheet()
            
            # Title
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                spaceAfter=30,
                textColor=colors.HexColor('#2E86AB')
            )
            
            story.append(Paragraph("Financial Planning & Analysis Report", title_style))
            
            # Company info
            company_info = st.session_state.get('report_config', {}).get('company_info', {})
            info_text = f"""
            <b>Company:</b> {company_info.get('company_name', '')}<br/>
            <b>Report Period:</b> {company_info.get('report_period', '')}<br/>
            <b>Prepared By:</b> {company_info.get('prepared_by', '')}<br/>
            <b>Date:</b> {company_info.get('prepared_date', '')}<br/>
            """
            
            story.append(Paragraph(info_text, styles["Normal"]))
            story.append(Spacer(1, 20))
            
            # Add sections
            if 'report_content' in st.session_state:
                content = st.session_state.report_content
                
                for section_key, section_content in content.items():
                    section_title = self.report_sections.get(section_key, {}).get('title', section_key)
                    
                    # Section title
                    story.append(Paragraph(section_title, styles['Heading2']))
                    story.append(Spacer(1, 12))
                    
                    # Section content
                    if section_content['type'] == 'text':
                        story.append(Paragraph(section_content['content'], styles['Normal']))
                        story.append(Spacer(1, 12))
                    
                    elif section_content['type'] == 'bullet_points':
                        for point in section_content['content']:
                            story.append(Paragraph(f"• {point}", styles['Normal']))
                        story.append(Spacer(1, 12))
                    
                    elif section_content['type'] == 'kpi_metrics':
                        # Add KPI table placeholder
                        kpi_data = [['KPI', 'Value', 'Change']]
                        if 'analysis_results' in st.session_state:
                            analysis = st.session_state.analysis_results
                            if 'kpi_dashboard' in analysis:
                                kpis = analysis['kpi_dashboard']
                                for kpi_name, kpi_info in list(kpis.items())[:5]:  # First 5 KPIs
                                    kpi_data.append([
                                        kpi_info.get('name', kpi_name),
                                        kpi_info.get('formatted_value', 'N/A'),
                                        ''
                                    ])
                        
                        kpi_table = Table(kpi_data, colWidths=[3*inch, 1.5*inch, 1.5*inch])
                        kpi_table.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2E86AB')),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 12),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                            ('GRID', (0, 0), (-1, -1), 1, colors.grey)
                        ]))
                        
                        story.append(kpi_table)
                        story.append(Spacer(1, 20))
            
            # Build PDF
            doc.build(story)
            
            # Get PDF bytes
            pdf_bytes = buffer.getvalue()
            buffer.close()
            
            return pdf_bytes
        
        except Exception as e:
            st.error(f"Error generating PDF: {str(e)}")
            return None
    
    def _generate_docx_report(self) -> Optional[bytes]:
        """Generate Word document report"""
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.style import WD_STYLE_TYPE
            import io
            
            # Create document
            doc = Document()
            
            # Add title
            title = doc.add_heading('Financial Planning & Analysis Report', 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Company info
            company_info = st.session_state.get('report_config', {}).get('company_info', {})
            doc.add_paragraph(f"Company: {company_info.get('company_name', '')}")
            doc.add_paragraph(f"Report Period: {company_info.get('report_period', '')}")
            doc.add_paragraph(f"Prepared By: {company_info.get('prepared_by', '')}")
            doc.add_paragraph(f"Date: {company_info.get('prepared_date', '')}")
            
            doc.add_paragraph()  # Empty paragraph
            
            # Add sections
            if 'report_content' in st.session_state:
                content = st.session_state.report_content
                
                for section_key, section_content in content.items():
                    section_title = self.report_sections.get(section_key, {}).get('title', section_key)
                    
                    # Add section title
                    doc.add_heading(section_title, level=1)
                    
                    # Add section content
                    if section_content['type'] == 'text':
                        doc.add_paragraph(section_content['content'])
                    
                    elif section_content['type'] == 'bullet_points':
                        for point in section_content['content']:
                            p = doc.add_paragraph(style='List Bullet')
                            p.add_run(point)
                    
                    elif section_content['type'] == 'kpi_metrics':
                        # Add KPI table
                        table = doc.add_table(rows=1, cols=3)
                        table.style = 'Light Grid Accent 1'
                        
                        # Header row
                        hdr_cells = table.rows[0].cells
                        hdr_cells[0].text = 'KPI'
                        hdr_cells[1].text = 'Value'
                        hdr_cells[2].text = 'Change'
                        
                        # Add KPI data
                        if 'analysis_results' in st.session_state:
                            analysis = st.session_state.analysis_results
                            if 'kpi_dashboard' in analysis:
                                kpis = analysis['kpi_dashboard']
                                for kpi_name, kpi_info in list(kpis.items())[:5]:
                                    row_cells = table.add_row().cells
                                    row_cells[0].text = kpi_info.get('name', kpi_name)
                                    row_cells[1].text = kpi_info.get('formatted_value', 'N/A')
                                    row_cells[2].text = ''
                        
                        doc.add_paragraph()
            
            # Save to buffer
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            
            return buffer.getvalue()
        
        except Exception as e:
            st.error(f"Error generating Word document: {str(e)}")
            return None
    
    def _generate_pptx_report(self) -> Optional[bytes]:
        """Generate PowerPoint presentation"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            import io
            
            # Create presentation
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "Financial Planning & Analysis Report"
            subtitle.text = "Executive Presentation"
            
            # Company info slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            title = slide.shapes.title
            title.text = "Report Information"
            
            company_info = st.session_state.get('report_config', {}).get('company_info', {})
            content = slide.placeholders[1]
            
            tf = content.text_frame
            tf.text = "Report Details:"
            
            p = tf.add_paragraph()
            p.text = f"Company: {company_info.get('company_name', '')}"
            
            p = tf.add_paragraph()
            p.text = f"Report Period: {company_info.get('report_period', '')}"
            
            p = tf.add_paragraph()
            p.text = f"Prepared By: {company_info.get('prepared_by', '')}"
            
            p = tf.add_paragraph()
            p.text = f"Date: {company_info.get('prepared_date', '')}"
            
            # Add content slides
            if 'report_content' in st.session_state:
                content_data = st.session_state.report_content
                
                for section_key, section_content in list(content_data.items())[:5]:  # First 5 sections
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    
                    section_title = self.report_sections.get(section_key, {}).get('title', section_key)
                    title = slide.shapes.title
                    title.text = section_title
                    
                    content = slide.placeholders[1]
                    tf = content.text_frame
                    
                    if section_content['type'] == 'text':
                        tf.text = section_content['content'][:500] + "..."  # Truncate for slide
                    
                    elif section_content['type'] == 'bullet_points':
                        for point in section_content['content'][:5]:  # First 5 points
                            p = tf.add_paragraph()
                            p.text = point
            
            # Summary slide
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            title.text = "Key Takeaways"
            
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = "Main Insights:"
            
            takeaways = [
                "Financial performance analyzed",
                "Key trends identified",
                "Recommendations provided",
                "Next steps outlined"
            ]
            
            for takeaway in takeaways:
                p = tf.add_paragraph()
                p.text = takeaway
            
            # Save to buffer
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            return buffer.getvalue()
        
        except Exception as e:
            st.error(f"Error generating PowerPoint: {str(e)}")
            return None
    
    def _generate_html_report(self) -> Optional[bytes]:
        """Generate HTML report"""
        try:
            import io
            
            # Get report data
            config = st.session_state.get('report_config', {})
            content = st.session_state.get('report_content', {})
            company_info = config.get('company_info', {})
            
            # Create HTML
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <title>FP&A Report - {company_info.get('report_period', '')}</title>
                <style>
                    body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
                    .header {{ background-color: #2E86AB; color: white; padding: 20px; border-radius: 5px; }}
                    .company-info {{ background-color: #f0f9ff; padding: 15px; border-left: 4px solid #2E86AB; margin: 20px 0; }}
                    .section {{ margin: 30px 0; }}
                    .section-title {{ color: #2E86AB; border-bottom: 2px solid #2E86AB; padding-bottom: 10px; }}
                    .kpi-table {{ width: 100%; border-collapse: collapse; margin: 15px 0; }}
                    .kpi-table th {{ background-color: #2E86AB; color: white; padding: 10px; text-align: left; }}
                    .kpi-table td {{ padding: 10px; border-bottom: 1px solid #ddd; }}
                    .bullet-point {{ margin: 5px 0; padding-left: 20px; }}
                    .footer {{ margin-top: 40px; padding-top: 20px; border-top: 1px solid #ddd; color: #666; font-size: 0.9em; }}
                </style>
            </head>
            <body>
                <div class="header">
                    <h1>Financial Planning & Analysis Report</h1>
                </div>
                
                <div class="company-info">
                    <p><strong>Company:</strong> {company_info.get('company_name', '')}</p>
                    <p><strong>Report Period:</strong> {company_info.get('report_period', '')}</p>
                    <p><strong>Prepared By:</strong> {company_info.get('prepared_by', '')}</p>
                    <p><strong>Date:</strong> {company_info.get('prepared_date', '')}</p>
                </div>
            """
            
            # Add sections
            for section_key, section_content in content.items():
                section_title = self.report_sections.get(section_key, {}).get('title', section_key)
                
                html_content += f"""
                <div class="section">
                    <h2 class="section-title">{section_title}</h2>
                """
                
                if section_content['type'] == 'text':
                    html_content += f"<p>{section_content['content']}</p>"
                
                elif section_content['type'] == 'bullet_points':
                    html_content += "<div>"
                    for point in section_content['content']:
                        html_content += f'<div class="bullet-point">• {point}</div>'
                    html_content += "</div>"
                
                elif section_content['type'] == 'kpi_metrics':
                    html_content += """
                    <table class="kpi-table">
                        <thead>
                            <tr>
                                <th>KPI</th>
                                <th>Value</th>
                                <th>Change</th>
                            </tr>
                        </thead>
                        <tbody>
                    """
                    
                    # Add KPI data
                    if 'analysis_results' in st.session_state:
                        analysis = st.session_state.analysis_results
                        if 'kpi_dashboard' in analysis:
                            kpis = analysis['kpi_dashboard']
                            for kpi_name, kpi_info in list(kpis.items())[:10]:  # First 10 KPIs
                                html_content += f"""
                                <tr>
                                    <td>{kpi_info.get('name', kpi_name)}</td>
                                    <td>{kpi_info.get('formatted_value', 'N/A')}</td>
                                    <td></td>
                                </tr>
                                """
                    
                    html_content += """
                        </tbody>
                    </table>
                    """
                
                html_content += "</div>"
            
            # Footer
            html_content += f"""
                <div class="footer">
                    <p>Generated by FP&A AI Agent on {datetime.now().strftime('%B %d, %Y %H:%M:%S')}</p>
                    <p>Confidential - Internal Use Only</p>
                </div>
            </body>
            </html>
            """
            
            # Convert to bytes
            buffer = io.BytesIO()
            buffer.write(html_content.encode('utf-8'))
            buffer.seek(0)
            
            return buffer.getvalue()
        
        except Exception as e:
            st.error(f"Error generating HTML: {str(e)}")
            return None
    
    def _generate_all_formats(self) -> Dict[str, bytes]:
        """Generate reports in all formats"""
        reports = {}
        
        # Generate PDF
        if reportlab_available:
            pdf_report = self._generate_pdf_report()
            if pdf_report:
                reports['pdf'] = pdf_report
        
        # Generate Word
        if python_docx_available:
            docx_report = self._generate_docx_report()
            if docx_report:
                reports['docx'] = docx_report
        
        # Generate PowerPoint
        if python_pptx_available:
            pptx_report = self._generate_pptx_report()
            if pptx_report:
                reports['pptx'] = pptx_report
        
        # Generate HTML
        html_report = self._generate_html_report()
        if html_report:
            reports['html'] = html_report
        
        return reports
    
    def _create_zip_file(self, reports: Dict[str, bytes], base_name: str) -> Optional[bytes]:
        """Create zip file from multiple reports"""
        try:
            import zipfile
            import io
            
            zip_buffer = io.BytesIO()
            
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                for format_name, report_data in reports.items():
                    filename = f"{base_name}.{format_name}"
                    zip_file.writestr(filename, report_data)
            
            zip_buffer.seek(0)
            return zip_buffer.getvalue()
        
        except Exception as e:
            st.error(f"Error creating zip file: {str(e)}")
            return None
    
    def _provide_download_link(self, file_bytes: bytes, filename: str, mime_type: str):
        """Provide download link for generated file"""
        b64 = base64.b64encode(file_bytes).decode()
        href = f'<a href="data:{mime_type};base64,{b64}" download="{filename}">Download {filename}</a>'
        st.markdown(href, unsafe_allow_html=True)
    
    def _report_library_tab(self) -> Dict[str, Any]:
        """
        Report library and templates
        """
        st.markdown("### 📁 Report Library")
        
        results = {}
        
        # Saved reports
        st.markdown("#### 💾 Saved Reports")
        
        # This would list saved reports in a real implementation
        st.info("Report library functionality coming soon. Generated reports will appear here.")
        
        # Report templates
        st.markdown("#### 📄 Report Templates")
        
        for template_key, template in self.report_templates.items():
            with st.expander(f"📋 {template['name']}"):
                st.write(f"**Description:** {template['description']}")
                st.write(f"**Target Audience:** {template['target_audience']}")
                st.write(f"**Sections:** {', '.join(template['sections'])}")
                
                if st.button(f"Use {template['name']}", key=f"use_{template_key}", use_container_width=True):
                    # Set this as current template
                    if 'report_config' not in st.session_state:
                        st.session_state.report_config = {}
                    st.session_state.report_config['template'] = template_key
                    st.session_state.report_config['sections'] = template['sections']
                    st.success(f"{template['name']} template loaded!")
        
        # Import/Export templates
        st.markdown("---")
        st.markdown("#### 📤 Import/Export")
        
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("📥 Import Template", use_container_width=True):
                st.info("Template import coming soon")
        
        with col2:
            if st.button("📤 Export Template", use_container_width=True):
                st.info("Template export coming soon")
        
        return results


# Streamlit integration function
def display_report_generator_section(dataframes: Dict[str, pd.DataFrame]) -> Dict[str, Any]:
    """
    Main function to display report generator interface
    """
    # Check for required data
    if not dataframes:
        st.warning("Please upload data first to generate reports.")
        return {}
    
    generator = ReportGenerator()
    report_results = generator.display_report_interface(dataframes)
    
    return report_results